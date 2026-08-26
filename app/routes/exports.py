"""Excel 匯入 / 匯出與週報投影片（PPTX）匯出。

匯出一律沿用列表當下的篩選條件（query string 相同），
所以「畫面上看到的就是匯出的內容」。
"""

import re
from datetime import date, datetime
from io import BytesIO

from fastapi import APIRouter, Depends, HTTPException, Query, Request
from fastapi.responses import HTMLResponse, Response
from fastapi.templating import Jinja2Templates
from loguru import logger
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from starlette.datastructures import UploadFile

from ..auth import require_scopes, require_user
from ..config import settings
from ..database import get_session
from ..fields import (
    FIELD_BY_KEY,
    LIST_COLUMNS,
    RFI_FIELDS,
    SLIDE_COLUMNS,
    diff_payload,
    display,
    short_week,
    to_iso_date,
    week_label,
)
from ..models import Rfi, RfiRevision
from ..query import filter_rows, parse_filters, sort_rows
from .rfis import (
    DASH_GROUP_LABELS,
    dashboard_stats,
    load_rows,
    next_rfi_no,
    sync_columns,
    _ctx,
)

router = APIRouter()

templates: Jinja2Templates = None  # type: ignore

# 匯出 / 匯入的欄位順序：編號、週別在前，其餘照欄位定義順序
EXPORT_KEYS = ["rfi_date", *[k for k in LIST_COLUMNS if k != "rfi_date"]]


async def _visible_rows(request: Request, session: AsyncSession):
    """套用與列表相同的篩選 / 排序，回傳要匯出的資料列。"""
    params = request.query_params
    rows = await load_rows(session)
    filters = parse_filters(params)
    rows = filter_rows(rows, filters, params.get("q", ""))
    return sort_rows(rows, params.get("sort", "week"), params.get("dir", "desc"))


def _stamp() -> str:
    return date.today().isoformat()


# ── Excel 匯出 ────────────────────────────────────────────────

def _harden(ws) -> None:
    """安全性：欄位值是使用者輸入，openpyxl 會把 "=" 開頭的字串存成公式，
    造成開啟匯出檔時執行公式（formula injection）。一律強制為純文字。"""
    for row in ws.iter_rows():
        for c in row:
            if c.data_type == "f":
                c.data_type = "s"


@router.get("/export/excel")
async def export_excel(
    request: Request,
    user: dict = Depends(require_user),
    session: AsyncSession = Depends(get_session),
):
    rows = await _visible_rows(request, session)
    if not rows:
        raise HTTPException(400, "目前的篩選條件下沒有資料可匯出")

    wb = Workbook()
    ws = wb.active
    ws.title = "RFI 列表"

    headers = ["RFI 編號", "週別"] + [FIELD_BY_KEY[k].label for k in EXPORT_KEYS] + [
        "附件", "最後更新", "更新者",
    ]
    ws.append(headers)
    for row in rows:
        r = row.obj
        ws.append(
            [r.rfi_no, r.week]
            + [row.v.get(k, "") for k in EXPORT_KEYS]
            + [
                "、".join(a.filename for a in r.attachments),
                r.updated_at.strftime("%Y-%m-%d %H:%M") if r.updated_at else "",
                r.updated_by,
            ]
        )

    head_font = Font(bold=True, color="FFFFFF")
    head_fill = PatternFill("solid", fgColor="1A73E8")
    for col in range(1, len(headers) + 1):
        c = ws.cell(row=1, column=col)
        c.font = head_font
        c.fill = head_fill
        c.alignment = Alignment(vertical="center", horizontal="center")
    for row in ws.iter_rows(min_row=2):
        for c in row:
            c.alignment = Alignment(vertical="top", wrap_text=True)

    widths = {"notes": 42, "risk": 24, "special": 18, "resolution": 14}
    ws.column_dimensions["A"].width = 14
    ws.column_dimensions["B"].width = 9
    for i, key in enumerate(EXPORT_KEYS, start=3):
        ws.column_dimensions[get_column_letter(i)].width = widths.get(key, 16)
    for i in range(len(EXPORT_KEYS) + 3, len(headers) + 1):
        ws.column_dimensions[get_column_letter(i)].width = 20
    ws.freeze_panes = "C2"
    _harden(ws)

    bio = BytesIO()
    wb.save(bio)
    logger.info("匯出 RFI Excel by {}（{} 筆）", user["sub"], len(rows))
    return _xlsx_response(bio.getvalue(), f"SA_RFI_List_{_stamp()}.xlsx")


def _xlsx_response(payload: bytes, filename: str) -> Response:
    return Response(
        payload,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.get("/dashboard/export")
async def export_dashboard(
    request: Request,
    user: dict = Depends(require_user),
    session: AsyncSession = Depends(get_session),
    group: str = Query("week"),
    year: str = Query(""),
):
    if group not in DASH_GROUP_LABELS:
        group = "week"
    stats = dashboard_stats(await load_rows(session), group, year)
    if not stats["table"]:
        raise HTTPException(400, "目前無統計數據可匯出")

    wb = Workbook()
    ws = wb.active
    ws.title = "統計看板"
    label = DASH_GROUP_LABELS[group]
    ws.append([label, "專案總數", "占比 (%)"])
    for r in stats["table"]:
        ws.append([r["category"], r["count"], f"{r['pct']}%"])
    ws.append(["合計", stats["total"], "100.0%"])

    head_font = Font(bold=True, color="FFFFFF")
    head_fill = PatternFill("solid", fgColor="1A73E8")
    for col in range(1, 4):
        c = ws.cell(row=1, column=col)
        c.font = head_font
        c.fill = head_fill
    ws.cell(row=ws.max_row, column=1).font = Font(bold=True)
    ws.column_dimensions["A"].width = 22
    ws.column_dimensions["B"].width = 14
    ws.column_dimensions["C"].width = 12
    _harden(ws)

    bio = BytesIO()
    wb.save(bio)
    suffix = f"_{year}" if year else ""
    logger.info("匯出 Dashboard Excel by {}（{}{}）", user["sub"], group, suffix)
    return _xlsx_response(bio.getvalue(), f"SA_RFI_Dashboard_{group}{suffix}_{_stamp()}.xlsx")


# ── 投影片（PPTX）匯出 ─────────────────────────────────────────

ROWS_PER_SLIDE = 6
TITLE_COLOR = RGBColor(0x00, 0x93, 0xDD)
HEADER_FILL = RGBColor(0x3F, 0x72, 0xAF)
LATIN_FONT = "Gill Sans MT"
EA_FONT = "Microsoft JhengHei"
# 各欄寬度（吋），總和 = 12.4，沿用 SA 原本週報的比例
COL_WIDTHS = [1.0, 0.9, 1.1, 0.9, 1.1, 0.6, 0.6, 0.8, 0.6, 1.2, 0.8, 2.8]


def _style_run(run, size: float, bold: bool, color: RGBColor) -> None:
    """設定字型（含中文的 East Asian 字型）與大小、顏色。"""
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.name = LATIN_FONT       # 產生 <a:latin>
    font.color.rgb = color

    # python-pptx 只管 latin 字型；中文要另外補 <a:ea>，否則 PowerPoint
    # 會用佈景主題的預設中文字型。<a:ea>/<a:cs> 依 schema 必須緊接在
    # <a:latin> 之後，所以用 addnext 依序插入而非直接 append。
    rpr = font._element          # CT_TextCharacterProperties
    prev = rpr.find(qn("a:latin"))
    if prev is None:
        return
    for tag, typeface in (("a:ea", EA_FONT), ("a:cs", LATIN_FONT)):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            prev.addnext(el)
        el.set("typeface", typeface)
        prev = el


def _fill_cell(cell, text: str, *, size: float, bold: bool, fill: RGBColor,
               color: RGBColor, align) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.04)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text or ""
    _style_run(run, size, bold, color)


def _week_range(rows) -> str:
    weeks = sorted({r.obj.week for r in rows if r.obj.week})
    if not weeks:
        return ""
    if len(weeks) == 1:
        return f" {short_week(weeks[0])}"
    return f" {short_week(weeks[0])}~{short_week(weeks[-1])}"


def build_deck(rows) -> bytes:
    """產生「Customer RFI Collection」週報投影片（16:9，每頁 6 筆）。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title = f"{settings.DECK_TITLE}{_week_range(rows)}"
    headers = [FIELD_BY_KEY[k].head for k in SLIDE_COLUMNS]
    notes_idx = SLIDE_COLUMNS.index("notes") if "notes" in SLIDE_COLUMNS else -1

    pages = [rows[i:i + ROWS_PER_SLIDE] for i in range(0, len(rows), ROWS_PER_SLIDE)]
    for page_no, page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(blank)

        box = slide.shapes.add_textbox(Inches(0.46), Inches(0.35), Inches(10.0), Inches(0.6))
        para = box.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = title
        _style_run(run, 32, True, TITLE_COLOR)

        table_shape = slide.shapes.add_table(
            len(page) + 1, len(headers),
            Inches(0.46), Inches(1.05), Inches(12.4), Inches(0.4 + 0.62 * len(page)),
        )
        table = table_shape.table
        # 關掉預設樣式的隔行網底，改由每格自行上色
        table.first_row = False
        table.horz_banding = False
        for i, width in enumerate(COL_WIDTHS[:len(headers)]):
            table.columns[i].width = Inches(width)
        table.rows[0].height = Inches(0.4)

        for col, text in enumerate(headers):
            _fill_cell(table.cell(0, col), text, size=12, bold=True,
                       fill=HEADER_FILL, color=RGBColor(0xFF, 0xFF, 0xFF),
                       align=PP_ALIGN.CENTER)

        for r, row in enumerate(page, start=1):
            table.rows[r].height = Inches(0.62)
            for col, key in enumerate(SLIDE_COLUMNS):
                _fill_cell(
                    table.cell(r, col), row.v.get(key, "") or "-",
                    size=10.5, bold=False,
                    fill=RGBColor(0xFF, 0xFF, 0xFF), color=RGBColor(0, 0, 0),
                    align=PP_ALIGN.LEFT if col == notes_idx else PP_ALIGN.CENTER,
                )

        if len(pages) > 1:
            foot = slide.shapes.add_textbox(
                Inches(11.6), Inches(6.95), Inches(1.3), Inches(0.35)
            )
            frun = foot.text_frame.paragraphs[0].add_run()
            frun.text = f"{page_no} / {len(pages)}"
            _style_run(frun, 11, False, RGBColor(0x5F, 0x63, 0x68))

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()


@router.get("/export/pptx")
async def export_pptx(
    request: Request,
    user: dict = Depends(require_user),
    session: AsyncSession = Depends(get_session),
):
    rows = await _visible_rows(request, session)
    if not rows:
        raise HTTPException(400, "目前的篩選條件下沒有資料可匯出")
    payload = build_deck(rows)
    logger.info("匯出 RFI 投影片 by {}（{} 筆 / {} 頁）",
                user["sub"], len(rows), -(-len(rows) // ROWS_PER_SLIDE))
    return Response(
        payload,
        media_type=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        headers={
            "Content-Disposition":
                f'attachment; filename="Customer_RFI_Collection_{_stamp()}.pptx"'
        },
    )


# ── Excel 匯入 ────────────────────────────────────────────────

# 表頭別名 → 欄位 key（相容 SA 舊有的 Excel 表頭與英文表頭）
_EXTRA_ALIASES: dict[str, list[str]] = {
    "rfi_date": ["日期", "date", "週別", "week", "rfi日期"],
    "client": ["終端客戶", "客戶", "client", "customer"],
    "vendor": ["面板廠", "vendor", "panel", "panelmaker"],
    "ic": ["ic型號", "ic", "icmodel", "型號"],
    "size": ["面板尺寸", "尺寸", "size", "panelsize"],
    "resolution": ["解析度", "resolution", "res"],
    "refresh_rate": ["頻率", "更新頻率", "hz", "frequency", "freq", "refreshrate"],
    "mux": ["mux"],
    "special": ["特殊", "特殊規格", "special"],
    "panel_type": ["type", "封裝", "封裝材質", "封裝/材質"],
    "product": ["產品", "產品類別", "product", "category"],
    "code": ["代號", "專案代號", "code", "project"],
    "notes": ["評估事項", "備註", "notes", "note", "remark"],
    "risk": ["風險", "風險評估", "risk"],
    "contact": ["客戶窗口", "窗口", "contact", "window"],
    "status": ["狀態", "處理狀態", "status"],
    "owner": ["負責sa", "負責人", "owner", "sa"],
}


def _norm(text) -> str:
    """正規化字串以便比對：去掉空白、底線、括號與各種分隔符號後轉小寫。

    這讓「COF 硬性」能對應到選項「COF，硬」、「IC 型號」能對應表頭「IC型號」。
    """
    return re.sub(r"[\s_()（）,，、/／·．.]+", "", str(text or "")).strip().lower()


def _header_map() -> dict[str, str]:
    mapping: dict[str, str] = {}
    for f in RFI_FIELDS:
        for alias in [f.key, f.label, f.head, *_EXTRA_ALIASES.get(f.key, [])]:
            mapping.setdefault(_norm(alias), f.key)
    return mapping


def _cell_text(value) -> str:
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.date().isoformat()
    if isinstance(value, date):
        return value.isoformat()
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def parse_import_row(raw: dict[str, str]) -> tuple[dict, list[str]]:
    """把一列 Excel 資料轉成 RFI data；回傳 (data, 修正說明)。"""
    warnings: list[str] = []
    data: dict = {}

    for f in RFI_FIELDS:
        text = raw.get(f.key, "")

        if f.type == "resolution":
            nums = re.findall(r"\d+", text)
            data[f.key] = {"w": nums[0] if nums else "", "h": nums[1] if len(nums) > 1 else ""}
            continue

        if f.key == "rfi_date":
            iso = to_iso_date(text)
            if text and not iso:
                warnings.append(f"日期「{text}」無法辨識，已留空")
            data[f.key] = iso
            continue

        if f.type == "number":
            match = re.search(r"\d+(?:\.\d+)?", text)
            data[f.key] = match.group(0) if match else ""
            continue

        if f.type == "select" and text and text not in f.options:
            fallback = _select_fallback(f.key, text)
            warnings.append(f"{f.label}「{text}」不在選項內，已改為「{fallback or '未填'}」")
            data[f.key] = fallback
            continue

        data[f.key] = text

    # 未填的選項欄位補上合理預設，讓匯入的資料可直接沿用
    for key, default in (("status", "評估中"), ("mux", "None")):
        if not data.get(key):
            data[key] = default
    return data, warnings


def _select_fallback(key: str, text: str) -> str:
    """匯入時遇到不在選項內的值，挑一個最接近的合法選項。"""
    options = FIELD_BY_KEY[key].options
    lowered = _norm(text)
    for opt in options:  # 大小寫 / 空白差異視為相同
        if _norm(opt) == lowered:
            return opt
    for opt in options:  # 部分包含（如「COF 柔性」→「COF，柔」）
        if _norm(opt) in lowered or lowered in _norm(opt):
            return opt
    return {"product": "其他", "status": "評估中", "mux": "None"}.get(key, "")


def _dedupe_key(data: dict) -> tuple:
    return (
        week_label(data.get("rfi_date", "")),
        data.get("client", "").lower(),
        data.get("vendor", "").lower(),
        data.get("ic", "").lower(),
        data.get("code", "").lower(),
    )


@router.get("/import", response_class=HTMLResponse)
async def import_form(
    request: Request, user: dict = Depends(require_scopes("write"))
):
    return templates.TemplateResponse(request, "import.html", _ctx(request, user))


@router.get("/import/template")
async def import_template(user: dict = Depends(require_user)):
    """下載匯入用的空白 Excel 範本（含一列範例）。"""
    wb = Workbook()
    ws = wb.active
    ws.title = "RFI 匯入範本"
    keys = EXPORT_KEYS
    ws.append([FIELD_BY_KEY[k].label for k in keys])
    sample = {
        "rfi_date": date.today().isoformat(), "client": "HP", "vendor": "BOE",
        "product": "OLED NB", "code": "Pro-14", "ic": "NT3670B", "size": "14.0",
        "resolution": "2560x1600", "refresh_rate": "60", "mux": "None",
        "panel_type": "COF，柔", "special": "HDR400", "status": "評估中",
        "owner": "Alex Lin", "contact": "Jenny Wu", "notes": "色域規格確認中。",
        "risk": "無",
    }
    ws.append([sample.get(k, "") for k in keys])
    for col, key in enumerate(keys, start=1):
        ws.column_dimensions[get_column_letter(col)].width = (
            42 if key == "notes" else 16
        )
        c = ws.cell(row=1, column=col)
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor="1A73E8")
    _harden(ws)

    bio = BytesIO()
    wb.save(bio)
    return _xlsx_response(bio.getvalue(), "SA_RFI_Import_Template.xlsx")


@router.post("/import", response_class=HTMLResponse)
async def import_excel(
    request: Request,
    user: dict = Depends(require_scopes("write")),
    session: AsyncSession = Depends(get_session),
):
    form = await request.form()
    upload = form.get("file")
    skip_dup = bool(form.get("skip_duplicates"))
    if not isinstance(upload, UploadFile) or not upload.filename:
        return templates.TemplateResponse(
            request, "import.html",
            _ctx(request, user, error="請先選擇要匯入的 Excel 檔案。"),
            status_code=400,
        )
    if not upload.filename.lower().endswith((".xlsx", ".xlsm")):
        return templates.TemplateResponse(
            request, "import.html",
            _ctx(request, user, error="僅支援 .xlsx / .xlsm 檔案。"),
            status_code=400,
        )

    content = await upload.read()
    if len(content) > settings.max_upload_bytes:
        return templates.TemplateResponse(
            request, "import.html",
            _ctx(request, user, error=f"檔案超過 {settings.MAX_UPLOAD_MB}MB 上限。"),
            status_code=400,
        )

    try:
        wb = load_workbook(BytesIO(content), data_only=True, read_only=True)
        ws = wb[wb.sheetnames[0]]
        sheet_rows = list(ws.iter_rows(values_only=True))
    except Exception as e:
        logger.warning("Excel 解析失敗 by {}：{}", user["sub"], e)
        return templates.TemplateResponse(
            request, "import.html",
            _ctx(request, user, error="Excel 解析失敗，請確認檔案未損毀且為標準格式。"),
            status_code=400,
        )

    if len(sheet_rows) < 2:
        return templates.TemplateResponse(
            request, "import.html",
            _ctx(request, user, error="檔案中沒有資料列（第 1 列須為表頭）。"),
            status_code=400,
        )

    header_map = _header_map()
    columns: dict[int, str] = {}
    for idx, cell in enumerate(sheet_rows[0]):
        key = header_map.get(_norm(cell))
        if key and key not in columns.values():
            columns[idx] = key
    if not columns:
        return templates.TemplateResponse(
            request, "import.html",
            _ctx(request, user,
                 error="無法辨識任何欄位表頭，請先下載範本或確認第 1 列為欄位名稱。"),
            status_code=400,
        )

    existing = {
        _dedupe_key(d or {}) for d in (await session.execute(select(Rfi.data))).scalars()
    }
    created, skipped, problems = 0, 0, []
    for line, raw_row in enumerate(sheet_rows[1:], start=2):
        raw = {columns[i]: _cell_text(v) for i, v in enumerate(raw_row) if i in columns}
        if not any(raw.values()):
            continue  # 空白列
        data, warnings = parse_import_row(raw)

        missing = [f.label for f in RFI_FIELDS if f.required and not display(f, data)]
        if missing:
            problems.append({"line": line, "level": "error",
                             "message": "缺少必填欄位：" + "、".join(missing)})
            continue

        key = _dedupe_key(data)
        if skip_dup and key in existing:
            skipped += 1
            problems.append({"line": line, "level": "skip",
                             "message": "與既有資料重複（同週別／客戶／面板廠／IC／代號），已略過"})
            continue
        existing.add(key)

        rfi = Rfi(created_by=user["sub"], updated_by=user["sub"], version=1)
        sync_columns(rfi, data)
        rfi.rfi_no = await next_rfi_no(session, rfi.week)
        session.add(rfi)
        await session.flush()
        session.add(RfiRevision(
            rfi_id=rfi.id, version=1, data=data, changes=diff_payload({}, data),
            action="import", note=f"由 Excel 匯入：{upload.filename}",
            edited_by=user["sub"],
        ))
        created += 1
        for w in warnings:
            problems.append({"line": line, "level": "warn", "message": w})

    await session.commit()
    logger.info("Excel 匯入 by {}：新增 {} 筆、略過 {} 筆、提示 {} 則",
                user["sub"], created, skipped, len(problems))
    return templates.TemplateResponse(
        request, "import.html",
        _ctx(request, user, created=created, skipped=skipped,
             problems=problems, filename=upload.filename, done=True),
    )
