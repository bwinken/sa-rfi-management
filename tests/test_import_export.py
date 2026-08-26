"""Excel 匯入 / 匯出與投影片匯出。

匯入的容錯規則（舊表頭、帶單位的值、選項容錯、去重）是 SA 從舊流程
搬過來的關鍵路徑，錯了會靜默地把資料吃壞，所以逐條釘住。
"""

from io import BytesIO

import pytest
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from sqlalchemy import select

from app.models import Rfi, RfiRevision
from tests.conftest import create_rfi

LEGACY_HEADERS = ["日期", "終端客戶", "面板廠", "IC型號", "面板尺寸", "解析度",
                  "頻率", "Mux", "特殊", "Type", "產品", "代號", "評估事項",
                  "風險", "客戶窗口"]


def _xlsx(rows: list[list], headers: list[str] = None) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.append(headers or LEGACY_HEADERS)
    for r in rows:
        ws.append(r)
    bio = BytesIO()
    wb.save(bio)
    return bio.getvalue()


async def _import(client, content: bytes, skip_dup: bool = True):
    data = {"skip_duplicates": "1"} if skip_dup else {}
    return await client.post(
        "/import",
        files={"file": ("legacy.xlsx", content,
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
        data=data,
    )


class TestImport:
    async def test_legacy_headers_recognised(self, client, session):
        content = _xlsx([["26W29", "HP", "BOE", "NT3670B", '14.0"', "2560x1600",
                          "60Hz", "None", "HDR400", "COF，柔", "OLED NB", "Pro-14",
                          "色域確認中", "無", "Alex"]])
        resp = await _import(client, content)
        assert resp.status_code == 200
        rfi = (await session.execute(select(Rfi))).scalar_one()
        assert rfi.client == "HP"
        assert rfi.week == "26W29"

    async def test_value_coercion(self, client, session):
        """帶單位、複合值要被拆解成乾淨的欄位值。"""
        content = _xlsx([["2026-08-24", "HP", "BOE", "NT1", '14.0"', "2560x1600",
                          "60Hz", "", "", "COF，柔", "OLED NB", "", "", "", ""]])
        await _import(client, content)
        rfi = (await session.execute(select(Rfi))).scalar_one()
        assert rfi.data["size"] == "14.0"
        assert rfi.data["refresh_rate"] == "60"
        assert rfi.data["resolution"] == {"w": "2560", "h": "1600"}

    async def test_select_value_fuzzy_matched(self, client, session):
        """『COF 硬性』要對應到選項『COF，硬』，並回報已修正。"""
        content = _xlsx([["26W30", "Dell", "AUO", "NT2", '27"', "3840x2160",
                          "240Hz", "", "", "COF 硬性", "OLED Monitor", "", "", "", ""]])
        resp = await _import(client, content)
        rfi = (await session.execute(select(Rfi))).scalar_one()
        assert rfi.data["panel_type"] == "COF，硬"
        assert "不在選項內" in resp.text

    async def test_unknown_product_falls_back(self, client, session):
        content = _xlsx([["26W30", "Sony", "JDI", "NT3", '6.1"', "1080x2340",
                          "120Hz", "", "", "COP，柔", "筆電", "", "", "", ""]])
        resp = await _import(client, content)
        rfi = (await session.execute(select(Rfi))).scalar_one()
        assert rfi.data["product"] == "其他"
        assert "已修正" in resp.text or "不在選項內" in resp.text

    async def test_row_missing_required_is_skipped_and_reported(self, client, session):
        content = _xlsx([
            ["26W31", "", "BOE", "NT9", '15"', "1920x1080", "60Hz", "", "",
             "COF，柔", "", "", "缺客戶與產品", "", ""],
        ])
        resp = await _import(client, content)
        assert (await session.execute(select(Rfi))).scalars().all() == []
        assert "缺少必填欄位" in resp.text

    async def test_blank_rows_ignored(self, client, session):
        content = _xlsx([
            ["26W29", "HP", "BOE", "NT1", '14"', "2560x1600", "60Hz", "", "",
             "COF，柔", "OLED NB", "", "", "", ""],
            ["", "", "", "", "", "", "", "", "", "", "", "", "", "", ""],
        ])
        await _import(client, content)
        assert len((await session.execute(select(Rfi))).scalars().all()) == 1

    async def test_duplicates_skipped(self, client, session):
        row = ["26W29", "HP", "BOE", "NT1", '14"', "2560x1600", "60Hz", "", "",
               "COF，柔", "OLED NB", "Pro-14", "", "", ""]
        content = _xlsx([row])
        await _import(client, content)
        resp = await _import(client, content)       # 同一份再匯一次
        assert len((await session.execute(select(Rfi))).scalars().all()) == 1
        assert "略過重複" in resp.text

    async def test_import_recorded_as_revision(self, client, session):
        content = _xlsx([["26W29", "HP", "BOE", "NT1", '14"', "2560x1600", "60Hz",
                          "", "", "COF，柔", "OLED NB", "", "", "", ""]])
        await _import(client, content)
        rev = (await session.execute(select(RfiRevision))).scalar_one()
        assert rev.action == "import"

    async def test_unrecognised_headers_rejected(self, client):
        content = _xlsx([["x"]], headers=["完全不相干的欄位"])
        resp = await _import(client, content)
        assert resp.status_code == 400
        assert "表頭" in resp.text

    async def test_non_xlsx_rejected(self, client):
        resp = await client.post(
            "/import", files={"file": ("evil.csv", b"a,b,c", "text/csv")},
        )
        assert resp.status_code == 400

    async def test_template_is_importable(self, client, session):
        """下載的範本應該能原樣匯入 —— 否則範本就是錯的。"""
        tpl = await client.get("/import/template")
        assert tpl.status_code == 200
        resp = await _import(client, tpl.content)
        assert "新增" in resp.text
        assert len((await session.execute(select(Rfi))).scalars().all()) == 1


class TestExcelExport:
    async def test_round_trip(self, client, rfi_form, session):
        """匯出的 Excel 再匯入回來，應該被判定為重複（欄位對得起來）。"""
        await create_rfi(client, rfi_form())
        exported = await client.get("/export/excel")
        assert exported.status_code == 200
        resp = await _import(client, exported.content)
        assert len((await session.execute(select(Rfi))).scalars().all()) == 1
        assert "略過重複" in resp.text

    async def test_respects_filters(self, client, rfi_form):
        await create_rfi(client, rfi_form(client="HP", vendor="BOE"))
        await create_rfi(client, rfi_form(client="Dell", vendor="AUO"))
        resp = await client.get("/export/excel?vendor=BOE")
        ws = load_workbook(BytesIO(resp.content)).active
        rows = list(ws.iter_rows(values_only=True))
        col = rows[0].index("終端客戶")      # 依表頭定位，欄序調整時測試不會假性失敗
        assert [r[col] for r in rows[1:]] == ["HP"]

    async def test_empty_result_is_400(self, client, rfi_form):
        await create_rfi(client, rfi_form(vendor="BOE"))
        resp = await client.get("/export/excel?vendor=NoSuchVendor")
        assert resp.status_code == 400

    async def test_formula_injection_neutralised(self, client, rfi_form, session):
        """使用者輸入若以 = 開頭，不可以在 Excel 裡變成公式。"""
        await create_rfi(client, rfi_form(notes="=1+1"))
        resp = await client.get("/export/excel")
        ws = load_workbook(BytesIO(resp.content)).active
        for row in ws.iter_rows():
            for cell in row:
                assert cell.data_type != "f"


class TestSlideExport:
    async def test_deck_structure(self, client, rfi_form):
        for i in range(7):     # 超過每頁 6 筆 → 應該分兩頁
            await create_rfi(client, rfi_form(client=f"Client{i}", code=f"C{i}"))
        resp = await client.get("/export/pptx")
        assert resp.status_code == 200
        prs = Presentation(BytesIO(resp.content))
        assert len(prs.slides) == 2
        assert round(prs.slide_width / 914400, 2) == 13.33     # 16:9

        table = next(s.table for s in prs.slides[0].shapes if s.has_table)
        assert len(table.rows) == 7          # 表頭 + 6 筆
        assert len(table.columns) == 12

    async def test_title_carries_week_range(self, client, rfi_form):
        await create_rfi(client, rfi_form(rfi_date="2026-08-10", code="A"))
        await create_rfi(client, rfi_form(rfi_date="2026-08-24", code="B"))
        resp = await client.get("/export/pptx")
        prs = Presentation(BytesIO(resp.content))
        texts = [sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame]
        assert any("W33~W35" in t for t in texts)

    async def test_chinese_font_is_declared(self, client, rfi_form):
        """python-pptx 只設 latin 字型，中文要另外寫 <a:ea>，否則會變成佈景預設字型。"""
        await create_rfi(client, rfi_form())
        resp = await client.get("/export/pptx")
        import zipfile
        xml = zipfile.ZipFile(BytesIO(resp.content)).read("ppt/slides/slide1.xml").decode()
        assert "Microsoft JhengHei" in xml

    async def test_empty_result_is_400(self, client):
        assert (await client.get("/export/pptx")).status_code == 400


class TestDashboardExport:
    @pytest.mark.parametrize("group", ["week", "product", "vendor", "client", "status", "owner"])
    async def test_all_groups(self, client, rfi_form, group):
        await create_rfi(client, rfi_form())
        resp = await client.get(f"/dashboard/export?group={group}")
        assert resp.status_code == 200

    async def test_invalid_group_falls_back(self, client, rfi_form):
        await create_rfi(client, rfi_form())
        assert (await client.get("/dashboard/export?group=bogus")).status_code == 200
